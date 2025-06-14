import os
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

class SlideGenerator:
    def __init__(self, output_dir="output/"):
        """Initialize the PowerPoint slide generator."""
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)  # Ensure output directory exists

    def create_slide_deck(self, title, bullet_points):
        """Create a PowerPoint presentation from summarized content."""
        prs = Presentation()
        
        # Title Slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Content Slides
        for i in range(0, len(bullet_points), 5):  # Limit 5 bullet points per slide
            slide_layout = prs.slide_layouts[1]  # Title + Content layout
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title_shape.text = f"{title} (Part {i//5 + 1})"
            for bullet in bullet_points[i:i+5]:  # Add max 5 points per slide
                p = content_shape.text_frame.add_paragraph()
                p.text = bullet

        # Save the presentation
        pptx_path = os.path.join(self.output_dir, "generated_presentation.pptx")
        prs.save(pptx_path)
        print(f"\n✅ PowerPoint slides saved to: {pptx_path}")

if __name__ == "__main__":
    # Example input
    sample_title = "AI-Based Slide Generation"
    sample_points = [
        "AI is revolutionizing industries through automation and data-driven decisions.",
        "Project goal: Automatic generation of PowerPoint slides from input documents.",
        "Supported file types: PDFs, DOCX, TXT, and CSV.",
        "Process: Text extraction and conversion into structured slides.",
        "Technology used: OpenAI's GPT models for text summarization and LlamaIndex."
    ]

    slide_generator = SlideGenerator()
    slide_generator.create_slide_deck(sample_title, sample_points)
